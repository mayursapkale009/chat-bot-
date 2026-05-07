#!/usr/bin/env python3
"""
Generate a professional PowerPoint presentation for:
"An Intelligent Multilingual Chatbot Platform Specifically Designed for Indian Languages"
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── colour palette ───────────────────────────────────────────────
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT_BLUE  = RGBColor(0x00, 0x96, 0xD6)   # vibrant cyan-blue
ACCENT_GOLD  = RGBColor(0xFF, 0xB8, 0x00)   # warm gold
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
SOFT_WHITE   = RGBColor(0xF0, 0xF0, 0xF5)
CARD_BG      = RGBColor(0x24, 0x24, 0x3E)   # slightly lighter card
GREEN        = RGBColor(0x00, 0xC8, 0x53)
RED_SOFT     = RGBColor(0xFF, 0x6B, 0x6B)
ORANGE       = RGBColor(0xFF, 0x9F, 0x43)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ─── helper functions ─────────────────────────────────────────────
def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, bullets, left, top, width, height,
                              font_size=16, color=WHITE, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_accent_bar(slide, top=Inches(1.55)):
    add_shape_rect(slide, Inches(0.8), top, Inches(1.8), Pt(4), ACCENT_BLUE)


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_section_label(slide, label):
    add_rounded_rect(slide, Inches(0.8), Inches(0.3), Inches(3.0), Inches(0.45),
                     ACCENT_BLUE)
    add_text_box(slide, Inches(0.9), Inches(0.32), Inches(2.8), Inches(0.4),
                 label, font_size=12, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)


def make_title(slide, title, subtitle=None):
    add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.8),
                 title, font_size=32, color=WHITE, bold=True, font_name="Calibri")
    add_accent_bar(slide, Inches(1.55))
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.75), Inches(11.5), Inches(0.5),
                     subtitle, font_size=16, color=LIGHT_GRAY)


def make_table(slide, rows, col_widths, left, top, header_color=ACCENT_BLUE):
    table_shape = slide.shapes.add_table(len(rows), len(col_widths), left, top,
                                          sum(col_widths), Inches(0.4) * len(rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Calibri"
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = WHITE
                paragraph.alignment = PP_ALIGN.CENTER

            # cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if ri == 0:
                cell_fill.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x48)
            else:
                cell_fill.fore_color.rgb = CARD_BG

            # borders
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
    return table_shape


TOTAL_SLIDES = 20

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)
# Decorative top bar
add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)
# Title
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.5), Inches(1.5),
             "An Intelligent Multilingual Chatbot Platform\nSpecifically Designed for Indian Languages",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Calibri")
# Gold accent line
add_shape_rect(slide, Inches(5.5), Inches(3.5), Inches(2.3), Pt(4), ACCENT_GOLD)
# Subtitle
add_text_box(slide, Inches(2), Inches(3.9), Inches(9.3), Inches(0.6),
             "IndicChat Pro — RAG + Ollama LLaMA3 + ChromaDB | Offline-First Multilingual AI",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
# Authors
add_text_box(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.5),
             "BE Final Year Project — Department of Computer Engineering",
             font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(5.6), Inches(9.3), Inches(0.5),
             "Guided by: Prof. [Guide Name]  |  Academic Year 2025-26",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
# Bottom bar
add_shape_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), ACCENT_GOLD)
add_slide_number(slide, 1, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Introduction
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "INTRODUCTION")
make_title(slide, "Why Multilingual Chatbots for India?")

bullets = [
    "🌍  India has 22 scheduled languages and 1,600+ dialects — most AI chatbots only serve English",
    "🚫  English-only systems create digital exclusion for 90%+ of India's population",
    "🔄  Code-switching (Hinglish) is the norm — existing NLP pipelines cannot handle mixed-language input",
    "📡  Cloud-dependent solutions fail in rural areas with poor connectivity",
    "🎯  Goal: Build an offline-first, multilingual chatbot supporting 11 Indian languages + Hinglish",
    "🧠  Key Technologies: RAG (Retrieval-Augmented Generation), Ollama LLaMA3, ChromaDB Vector Store",
    "📊  Achieves 85.2% QA accuracy — a 24.5% improvement over baseline systems",
]
add_bullet_slide_content(slide, bullets, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
                          font_size=17)
add_slide_number(slide, 2, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — Background / Problem Statement
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "BACKGROUND & PROBLEM")
make_title(slide, "Problem Statement & Motivation")

# Left card
card = add_rounded_rect(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), CARD_BG, ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(2.35), Inches(5.1), Inches(0.4),
             "Key Problems", font_size=20, color=ACCENT_GOLD, bold=True)
problems = [
    "❌ Lack of Multilingual Support: Most chatbots are English-only",
    "❌ Poor Context Awareness: Systems fail across language/cultural contexts",
    "❌ Scalability Issues: Frameworks struggle with language diversity",
    "❌ English Bias in LLMs: Models internally use English representations",
    "❌ ASR Failures: High WER for Indic languages (179.5% for Kannada)",
    "❌ No Offline Capability: Cloud dependency limits rural deployment",
]
add_bullet_slide_content(slide, problems, Inches(1.0), Inches(2.85), Inches(5.1), Inches(3.5),
                          font_size=14, color=SOFT_WHITE)

# Right card
card2 = add_rounded_rect(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.5), CARD_BG, ACCENT_GOLD)
add_text_box(slide, Inches(7.0), Inches(2.35), Inches(5.1), Inches(0.4),
             "Technical Challenges", font_size=20, color=ACCENT_GOLD, bold=True)
challenges = [
    "⚡ Rich morphology & varied syntax in Indian languages",
    "⚡ Code-switching (Hinglish) detection & processing",
    "⚡ Limited training datasets for low-resource languages",
    "⚡ Tokenization issues with Indic scripts",
    "⚡ Spelling variations across regional dialects",
    "⚡ Need for cultural sensitivity & fairness",
]
add_bullet_slide_content(slide, challenges, Inches(7.0), Inches(2.85), Inches(5.1), Inches(3.5),
                          font_size=14, color=SOFT_WHITE)
add_slide_number(slide, 3, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — Background Theory / Key Technologies
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "BACKGROUND THEORY")
make_title(slide, "Core Technologies & Concepts")

techs = [
    ("RAG", "Retrieval-Augmented Generation — combines retrieval from vector DB with LLM generation to reduce hallucinations"),
    ("Ollama LLaMA3", "Open-source LLM running locally via Ollama — enables offline, air-gapped deployment"),
    ("ChromaDB", "Lightweight vector database for storing & querying document embeddings locally"),
    ("MuRIL BERT", "Multilingual Representations for Indian Languages — pre-trained on 17 Indian languages"),
    ("LangDetect", "Language identification library for detecting input language including Hinglish code-mixing"),
    ("Whisper ASR", "OpenAI's speech recognition model — fine-tuned for Indic languages to reduce WER"),
]

y_start = 2.2
for i, (title, desc) in enumerate(techs):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(y_start + row * 1.55)
    add_rounded_rect(slide, x, y, Inches(5.8), Inches(1.35), CARD_BG, ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.35),
                 title, font_size=18, color=ACCENT_GOLD, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.4), Inches(0.75),
                 desc, font_size=13, color=SOFT_WHITE)

add_slide_number(slide, 4, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Literature Survey 1
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "LITERATURE SURVEY")
make_title(slide, "Related Works — Multilingual Chatbot Systems")

rows = [
    ["Source", "System / Model", "Key Contribution", "Limitation"],
    ["[1]", "MuRIL BERT FAQ", "60.73% Top-1 accuracy on multilingual FAQ retrieval", "No code-switching support"],
    ["[3]", "ChatSense", "Dynamic knowledge graphs + Gemini for voice interaction", "Cloud-dependent"],
    ["[4]", "RAG + Knowledge Graph", "Structured domain knowledge improves reasoning", "Limited language support"],
    ["[7]", "Multi-Agent Hybrid", "Combines generative + retrieval methods", "Complex architecture overhead"],
    ["[11]", "IndicBERT / IndicNLPSuite", "Compact 12M-18M model, IndicGLUE benchmark", "No chatbot system built"],
    ["[12]", "ChatGPT + Vector DB", "LangChain + FAISS/Pinecone for domain QA", "Proprietary, not offline"],
]
make_table(slide, rows,
           [Inches(1.0), Inches(2.5), Inches(5.0), Inches(3.0)],
           Inches(0.8), Inches(2.2))
add_slide_number(slide, 5, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Literature Survey 2
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "LITERATURE SURVEY")
make_title(slide, "Related Works — ASR & NLP for Indic Languages")

rows2 = [
    ["Source", "Focus Area", "Key Findings"],
    ["[2]", "Hindi NLP Chatbots", "NLP-based chatbots for healthcare & education in Hindi"],
    ["[5]", "LLM Bias Analysis", "English-centric representations degrade Indic performance"],
    ["[6]", "Whisper ASR Fine-tuning", "Joint multilingual training reduces WER by 40-60%"],
    ["[8]", "Cross-lingual Transfer", "Adapter-based fine-tuning enables low-resource language support"],
    ["[9]", "Chatbot Evaluation", "Turing Test frameworks for conversational AI assessment"],
    ["[13]", "Survey & Trends", "K-means clustering & ontology analysis of chatbot research"],
]
make_table(slide, rows2,
           [Inches(1.0), Inches(3.5), Inches(7.0)],
           Inches(0.8), Inches(2.2))

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
             "Research Gap: No existing system combines offline LLM + RAG + 11 Indian languages + Hinglish code-mixing in a single platform.",
             font_size=16, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 6, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Proposed Model Overview
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "PROPOSED MODEL")
make_title(slide, "IndicChat Pro — System Overview")

features = [
    ("🌐", "11 Languages + Hinglish", "Hindi, Bengali, Tamil, Telugu, Marathi,\nGujarati, Kannada, Malayalam, Punjabi,\nOdia, Assamese + Hinglish code-mix"),
    ("🧠", "RAG Architecture", "Retrieval-Augmented Generation\nwith ChromaDB vector store for\naccurate, context-aware responses"),
    ("💻", "Offline-First Design", "Ollama LLaMA3 runs locally\nAir-gapped deployment ready\n~1.8s response latency"),
    ("📊", "85.2% Accuracy", "Outperforms MuRIL (60.73%) and\nIndicBERT baselines on QA tasks\n92% Hinglish handling rate"),
]

for i, (icon, title, desc) in enumerate(features):
    x = Inches(0.6 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(2.3), Inches(2.9), Inches(4.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.3), Inches(2.5), Inches(2.3), Inches(0.5),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.1), Inches(2.5), Inches(0.5),
                 title, font_size=16, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(2.5),
                 desc, font_size=13, color=SOFT_WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 7, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — System Architecture (Fig 1)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "PROPOSED MODEL")
make_title(slide, "System Architecture — IndicChat Pro", "Fig 1: End-to-End Architecture with RAG + Ollama LLaMA3 + ChromaDB")

# Architecture flow boxes
flow_items = [
    ("User Input\n(11 Languages\n+ Hinglish)", ACCENT_BLUE),
    ("Language\nDetection\n(LangDetect)", RGBColor(0x6C, 0x5C, 0xE7)),
    ("NLU Module\n(MuRIL BERT\nEmbeddings)", RGBColor(0x00, 0xB8, 0x94)),
    ("Vector Search\n(ChromaDB\nCosine Sim.)", ACCENT_GOLD),
    ("LLM Generation\n(Ollama\nLLaMA3)", RGBColor(0xE1, 0x72, 0x55)),
    ("Response\n(Multilingual\nOutput)", GREEN),
]

y_top = Inches(2.8)
box_w = Inches(1.7)
box_h = Inches(1.6)
gap = Inches(0.25)
start_x = Inches(0.6)

for i, (label, color) in enumerate(flow_items):
    x = start_x + i * (box_w + gap)
    add_rounded_rect(slide, x, y_top, box_w, box_h, color)
    add_text_box(slide, x + Inches(0.05), y_top + Inches(0.2), box_w - Inches(0.1), box_h - Inches(0.2),
                 label, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < len(flow_items) - 1:
        arrow_x = x + box_w
        add_text_box(slide, arrow_x, y_top + Inches(0.55), gap, Inches(0.4),
                     "→", font_size=24, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# Bottom knowledge layer
add_rounded_rect(slide, Inches(0.6), Inches(5.0), Inches(11.7), Inches(1.8), CARD_BG, ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.1), Inches(3.0), Inches(0.35),
             "Knowledge Layer Pipeline:", font_size=16, color=ACCENT_GOLD, bold=True)

kl_items = ["Business\nDocs (PDF)", "Text\nExtraction", "Chunking\n(512 tokens)", "Embedding\n(MuRIL)", "ChromaDB\nVector Store"]
for i, item in enumerate(kl_items):
    x = Inches(0.8 + i * 2.3)
    add_rounded_rect(slide, x, Inches(5.55), Inches(1.8), Inches(1.0),
                     RGBColor(0x34, 0x34, 0x55), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.05), Inches(5.6), Inches(1.7), Inches(0.9),
                 item, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(kl_items) - 1:
        add_text_box(slide, x + Inches(1.8), Inches(5.8), Inches(0.5), Inches(0.4),
                     "→", font_size=20, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 8, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Multilingual NLU Module (Fig 2)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "PROPOSED MODEL")
make_title(slide, "Multilingual NLU Module", "Fig 2: Language Detection & Hinglish Processing Pipeline")

# NLU Flow
nlu_steps = [
    ("User Query\n(Any Language)", ACCENT_BLUE, "Input can be Hindi,\nTamil, Hinglish, etc."),
    ("LangDetect\nModule", RGBColor(0x6C, 0x5C, 0xE7), "Identifies language using\nstatistical analysis + regex\nfor code-mixed text"),
    ("Hinglish\nDetector", ACCENT_GOLD, "Regex-based detection of\ncode-switching patterns\n92% accuracy"),
    ("MuRIL BERT\nEncoder", RGBColor(0x00, 0xB8, 0x94), "Generates 256-dim\nvector embeddings\nfor semantic matching"),
    ("Intent\nClassifier", RGBColor(0xE1, 0x72, 0x55), "Maps query to\npredefined categories\nfor retrieval routing"),
]

for i, (title, color, desc) in enumerate(nlu_steps):
    x = Inches(0.5 + i * 2.55)
    add_rounded_rect(slide, x, Inches(2.3), Inches(2.25), Inches(1.2), color)
    add_text_box(slide, x + Inches(0.1), Inches(2.4), Inches(2.05), Inches(1.0),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(nlu_steps) - 1:
        add_text_box(slide, x + Inches(2.25), Inches(2.7), Inches(0.3), Inches(0.3),
                     "→", font_size=20, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
    # Description card below
    add_rounded_rect(slide, x, Inches(3.7), Inches(2.25), Inches(1.5), CARD_BG)
    add_text_box(slide, x + Inches(0.1), Inches(3.8), Inches(2.05), Inches(1.3),
                 desc, font_size=12, color=SOFT_WHITE, alignment=PP_ALIGN.CENTER)

# Supported languages bar
add_rounded_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.3), CARD_BG, ACCENT_GOLD)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(4.0), Inches(0.35),
             "Supported Languages (11 + Hinglish):", font_size=15, color=ACCENT_GOLD, bold=True)
langs = "Hindi • Bengali • Tamil • Telugu • Marathi • Gujarati • Kannada • Malayalam • Punjabi • Odia • Assamese • Hinglish (Code-Mixed)"
add_text_box(slide, Inches(0.7), Inches(6.15), Inches(11.8), Inches(0.5),
             langs, font_size=14, color=WHITE, alignment=PP_ALIGN.LEFT)
add_slide_number(slide, 9, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Knowledge Layer (Fig 3)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "PROPOSED MODEL")
make_title(slide, "Knowledge Layer Pipeline", "Fig 3: Business Documents → Vector Store Representations")

# Step cards
steps = [
    ("Step 1", "Document Ingestion", "Business PDFs, FAQs, and\nunstructured documents are\ncollected as knowledge sources", ACCENT_BLUE),
    ("Step 2", "Text Extraction", "PyPDF2 / OCR extracts raw\ntext from documents,\nhandling multiple scripts", RGBColor(0x6C, 0x5C, 0xE7)),
    ("Step 3", "Chunking", "Text is split into 512-token\nchunks with overlap for\ncontext preservation", RGBColor(0x00, 0xB8, 0x94)),
    ("Step 4", "Embedding", "MuRIL BERT generates\n256-dimensional vector\nembeddings for each chunk", ACCENT_GOLD),
    ("Step 5", "Vector Storage", "Embeddings stored in\nChromaDB with metadata\nfor fast cosine similarity search", RGBColor(0xE1, 0x72, 0x55)),
]

for i, (step, title, desc, color) in enumerate(steps):
    x = Inches(0.4 + i * 2.55)
    # Step number
    add_rounded_rect(slide, x, Inches(2.3), Inches(2.3), Inches(0.45), color)
    add_text_box(slide, x + Inches(0.1), Inches(2.32), Inches(2.1), Inches(0.4),
                 f"{step}: {title}", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_rounded_rect(slide, x, Inches(2.9), Inches(2.3), Inches(1.8), CARD_BG, color)
    add_text_box(slide, x + Inches(0.1), Inches(3.0), Inches(2.1), Inches(1.6),
                 desc, font_size=12, color=SOFT_WHITE, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(2.3), Inches(3.5), Inches(0.25), Inches(0.4),
                     "→", font_size=18, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# Query flow at bottom
add_rounded_rect(slide, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.8), CARD_BG, ACCENT_BLUE)
add_text_box(slide, Inches(0.6), Inches(5.3), Inches(4.0), Inches(0.35),
             "Runtime Query Flow:", font_size=16, color=ACCENT_GOLD, bold=True)

query_flow = [
    "User Query → MuRIL Embedding → Cosine Similarity Search in ChromaDB → Top-K Relevant Chunks → LLaMA3 Generates Answer"
]
add_text_box(slide, Inches(0.6), Inches(5.75), Inches(12.0), Inches(0.4),
             query_flow[0], font_size=15, color=WHITE)

math_text = "Similarity = cos(θ) = (A·B) / (||A|| × ||B||)    where A = query embedding, B = document embedding"
add_text_box(slide, Inches(0.6), Inches(6.25), Inches(12.0), Inches(0.4),
             math_text, font_size=14, color=LIGHT_GRAY)
add_slide_number(slide, 10, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — RAG Architecture Detail
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "PROPOSED MODEL")
make_title(slide, "RAG Architecture — Retrieval-Augmented Generation")

# Left: How RAG works
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.8), Inches(4.5), CARD_BG, ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(2.4), Inches(5.4), Inches(0.4),
             "How RAG Works in IndicChat Pro", font_size=18, color=ACCENT_GOLD, bold=True)

rag_steps = [
    "1️⃣  User submits query in any supported Indian language",
    "2️⃣  LangDetect identifies language; Hinglish is parsed via regex",
    "3️⃣  MuRIL BERT encodes query into 256-dim vector",
    "4️⃣  ChromaDB performs cosine similarity search (Top-K=5)",
    "5️⃣  Retrieved context chunks are injected into LLM prompt",
    "6️⃣  Ollama LLaMA3 generates grounded, factual response",
    "7️⃣  Response delivered in user's detected language",
]
add_bullet_slide_content(slide, rag_steps, Inches(1.0), Inches(2.95), Inches(5.4), Inches(3.5),
                          font_size=14, color=SOFT_WHITE)

# Right: Benefits
add_rounded_rect(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5), CARD_BG, ACCENT_GOLD)
add_text_box(slide, Inches(7.2), Inches(2.4), Inches(5.1), Inches(0.4),
             "Why RAG Over Pure LLM?", font_size=18, color=ACCENT_GOLD, bold=True)

benefits = [
    "✅ Reduces hallucinations by grounding in real documents",
    "✅ Domain-specific without expensive fine-tuning",
    "✅ Easy knowledge updates — just add new documents",
    "✅ Works offline with local vector DB (ChromaDB)",
    "✅ Transparent — retrieved sources can be cited",
    "✅ Lower compute than full model re-training",
]
add_bullet_slide_content(slide, benefits, Inches(7.2), Inches(2.95), Inches(5.1), Inches(3.5),
                          font_size=14, color=SOFT_WHITE)
add_slide_number(slide, 11, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — Model Selection & MuRIL BERT
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "PROPOSED MODEL")
make_title(slide, "Model Selection — Why MuRIL BERT?")

rows_model = [
    ["Aspect", "MuRIL BERT", "XLM-RoBERTa", "IndicBERT", "Monsoon-NLP"],
    ["Architecture", "BERT-base", "RoBERTa-large", "ALBERT-based", "BERT-based"],
    ["Languages", "17 Indian + English", "~100 languages", "12 Indian", "Hindi-focused"],
    ["Parameters", "~110M", "~125M", "12M-18M", "~110M"],
    ["Training Data", "Mono + Trans + Translit", "2.5TB Common Crawl", "IndicCorp", "Hindi corpus"],
    ["Top-1 Accuracy", "60.73%", "39.75%", "39.75%", "14.63%"],
    ["Indic Optimized", "✅ Yes", "❌ No", "✅ Yes", "⚠️ Partial"],
    ["Selected?", "✅ CHOSEN", "—", "—", "—"],
]
make_table(slide, rows_model,
           [Inches(2.0), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.0)],
           Inches(0.8), Inches(2.2))

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
             "MuRIL BERT was selected for its superior performance on Indian language tasks, trained on monolingual, translated, and transliterated data for 17 Indian languages.",
             font_size=14, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 12, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — Offline LLM & Deployment
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "PROPOSED MODEL")
make_title(slide, "Offline-First Deployment with Ollama LLaMA3")

# Three column cards
cards_data = [
    ("Ollama Runtime", ACCENT_BLUE,
     ["Runs LLaMA3 locally", "No cloud API needed", "Air-gapped compatible",
      "Single binary deployment", "GPU acceleration support"]),
    ("Privacy & Security", RGBColor(0x00, 0xB8, 0x94),
     ["All data stays on-device", "No external API calls", "Compliant with data laws",
      "Suitable for government", "Enterprise-ready"]),
    ("Performance", ACCENT_GOLD,
     ["~1.8s response latency", "Handles concurrent users", "Low memory footprint",
      "Efficient token generation", "Scalable horizontally"]),
]

for i, (title, color, items) in enumerate(cards_data):
    x = Inches(0.6 + i * 4.1)
    add_rounded_rect(slide, x, Inches(2.3), Inches(3.8), Inches(4.2), CARD_BG, color)
    add_text_box(slide, x + Inches(0.2), Inches(2.45), Inches(3.4), Inches(0.4),
                 title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape_rect(slide, x + Inches(0.5), Inches(2.95), Inches(2.8), Pt(2), color)
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.3), Inches(3.15 + j * 0.55), Inches(3.2), Inches(0.4),
                     f"• {item}", font_size=14, color=SOFT_WHITE)

add_slide_number(slide, 13, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — Results: Performance Comparison Table
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "RESULTS")
make_title(slide, "Performance: IndicChat Pro vs. Baselines", "Table IV: Comprehensive Performance Comparison")

rows_perf = [
    ["Component", "Metric", "IndicChat Pro", "Baseline"],
    ["Accuracy", "RAG-based QA", "85.2%", "60.7%"],
    ["Language Support", "Multilingual", "11 langs + Hinglish (92%)", "9 languages"],
    ["Latency", "Response Time", "~1.8 seconds (offline)", "Cloud-based (higher)"],
    ["Benchmark", "Dataset", "IndicGLUE", "IndicGLUE"],
    ["EM Score", "Exact Match", "0.85", "0.61"],
    ["F1 Score", "F1", "0.85", "0.61"],
]
make_table(slide, rows_perf,
           [Inches(2.5), Inches(2.5), Inches(3.5), Inches(3.0)],
           Inches(0.8), Inches(2.2))

# Highlight box
add_rounded_rect(slide, Inches(3.5), Inches(6.0), Inches(6.3), Inches(0.8), GREEN)
add_text_box(slide, Inches(3.7), Inches(6.05), Inches(5.9), Inches(0.7),
             "🏆 24.5% improvement in QA accuracy over baseline | 92% Hinglish handling",
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 14, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — Results: FAQ Chatbot Performance
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "RESULTS")
make_title(slide, "FAQ Chatbot — Top-K Accuracy Comparison", "Table VI: Performance Across Models")

rows_faq = [
    ["Model", "Top-1 Accuracy", "Top-2 Accuracy", "Top-3 Accuracy"],
    ["MuRIL BERT", "60.73%", "70.73%", "76.34%"],
    ["XLM-RoBERTa", "39.75%", "46.34%", "49.02%"],
    ["IndicBERT", "39.75%", "46.33%", "48.72%"],
    ["Monsoon-NLP", "14.63%", "15.60%", "16.09%"],
]
make_table(slide, rows_faq,
           [Inches(3.0), Inches(2.8), Inches(2.8), Inches(2.8)],
           Inches(0.8), Inches(2.2))

# Visual bar chart representation using shapes
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.0), Inches(0.4),
             "Top-1 Accuracy — Visual Comparison:", font_size=16, color=ACCENT_GOLD, bold=True)

bars = [("MuRIL", 60.73, GREEN), ("XLM-R", 39.75, ACCENT_BLUE),
        ("IndicBERT", 39.75, RGBColor(0x6C, 0x5C, 0xE7)), ("Monsoon", 14.63, RED_SOFT)]
for i, (name, val, color) in enumerate(bars):
    y = Inches(5.0 + i * 0.52)
    add_text_box(slide, Inches(0.8), y, Inches(1.5), Inches(0.4), name, font_size=13, color=WHITE)
    bar_width = Inches(val / 100 * 8.0)
    add_rounded_rect(slide, Inches(2.5), y + Inches(0.05), bar_width, Inches(0.3), color)
    add_text_box(slide, Inches(2.5) + bar_width + Inches(0.1), y, Inches(1.0), Inches(0.4),
                 f"{val}%", font_size=13, color=color, bold=True)

add_slide_number(slide, 15, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 16 — Results: ASR Performance (Whisper)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "RESULTS")
make_title(slide, "ASR Performance — Whisper Fine-Tuning for Indic Languages", "Table VII: Word Error Rate (WER) Comparison — Lower is Better")

rows_asr = [
    ["System", "English (En)", "Hindi (Hi)", "Bengali (Bn)", "Telugu (Te)", "Kannada (Kn)"],
    ["Whisper (Baseline)", "49.6", "62.8", "116.8", "114.9", "179.5"],
    ["Individual Fine-Tune", "—", "36.6", "55.1", "68.2", "63.5"],
    ["Proposed v1.1 (Unified)", "32.5", "40.0", "55.9", "85.6", "73.0"],
]
make_table(slide, rows_asr,
           [Inches(2.8), Inches(2.0), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)],
           Inches(0.8), Inches(2.2))

# Language ID table
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(6.0), Inches(0.4),
             "Table VIII: Language Identification — Diarization Error Rate (DER)", font_size=15, color=ACCENT_GOLD, bold=True)

rows_der = [
    ["System", "DER (Lower = Better)"],
    ["Whisper (Baseline)", "24.04"],
    ["Proposed v1 (Joint ASR/LID)", "18.79"],
    ["Proposed v2 (0.8 LID / 0.2 ASR)", "18.03"],
    ["LID-only Training", "23.46"],
]
make_table(slide, rows_der,
           [Inches(5.5), Inches(3.5)],
           Inches(0.8), Inches(5.0))
add_slide_number(slide, 16, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 17 — Results: Cross-System Comparison
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "RESULTS")
make_title(slide, "Cross-System Comparison — IndicChat Pro vs Research Baselines", "Table V: Comprehensive Feature & Performance Matrix")

rows_cross = [
    ["Component", "IndicChat Pro", "MuRIL FAQ [1]", "IndicNLPSuite [11]"],
    ["Datasets", "MCA FAQs + Business QA\n+ IndicGLUE", "Ministry of Corporate\nAffairs FAQs", "IndicGLUE, IndicCorp"],
    ["Languages", "11 Indic + Hinglish\n(92% code-mix)", "9 Indian languages", "11 major Indian\n+ English"],
    ["Accuracy", "85.2% QA accuracy", "60.73% Top-1 FAQ", "Competitive NLU\n(no chatbot focus)"],
    ["Base Model", "RAG + Ollama LLaMA3\n+ ChromaDB (offline)", "MuRIL BERT\n(extractive QA)", "IndicBERT\n(ALBERT)"],
    ["Code-Mixing", "Full support\n(regex + langdetect)", "No support\n(future goal)", "Monolingual only"],
    ["Latency", "1.8s offline,\nair-gapped ready", "Reduced NMT\noverhead", "Compact\n(12M-18M params)"],
]
make_table(slide, rows_cross,
           [Inches(2.2), Inches(3.2), Inches(3.0), Inches(3.0)],
           Inches(0.8), Inches(2.2))
add_slide_number(slide, 17, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 18 — Advantages & Disadvantages
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "COMPARISON")
make_title(slide, "Advantages & Disadvantages")

# Left: Advantages
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), CARD_BG, GREEN)
add_text_box(slide, Inches(1.0), Inches(2.4), Inches(5.1), Inches(0.4),
             "✅ Advantages", font_size=22, color=GREEN, bold=True)

advantages = [
    "✅ 85.2% QA accuracy — 24.5% improvement over baselines",
    "✅ Supports 11 Indian languages + Hinglish code-mixing",
    "✅ Fully offline — no cloud dependency, works in rural areas",
    "✅ Privacy-preserving — all data stays on-device",
    "✅ RAG reduces hallucinations with grounded responses",
    "✅ Easy knowledge updates without model retraining",
    "✅ Air-gapped deployment for government/enterprise use",
]
add_bullet_slide_content(slide, advantages, Inches(1.0), Inches(2.95), Inches(5.1), Inches(3.5),
                          font_size=14, color=SOFT_WHITE)

# Right: Disadvantages
add_rounded_rect(slide, Inches(6.8), Inches(2.3), Inches(5.5), Inches(4.5), CARD_BG, RED_SOFT)
add_text_box(slide, Inches(7.0), Inches(2.4), Inches(5.1), Inches(0.4),
             "⚠️ Limitations", font_size=22, color=RED_SOFT, bold=True)

disadvantages = [
    "⚠️ Limited training data for some low-resource languages",
    "⚠️ Local LLM requires capable hardware (GPU recommended)",
    "⚠️ Spelling variations across dialects still challenging",
    "⚠️ No voice interaction yet (future ASR integration)",
    "⚠️ Cultural nuances may be missed in some contexts",
    "⚠️ Code-switching accuracy can degrade with more languages",
    "⚠️ Model size constraints limit response sophistication",
]
add_bullet_slide_content(slide, disadvantages, Inches(7.0), Inches(2.95), Inches(5.1), Inches(3.5),
                          font_size=14, color=SOFT_WHITE)
add_slide_number(slide, 18, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 19 — Conclusion & Future Work
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "CONCLUSION & FUTURE WORK")
make_title(slide, "Conclusion & Future Directions")

# Left: Conclusion
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), CARD_BG, ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(2.4), Inches(5.1), Inches(0.4),
             "🎯 Key Conclusions", font_size=20, color=ACCENT_BLUE, bold=True)

conclusions = [
    "• Multilingual chatbots are critical for India's digital inclusion",
    "• RAG + offline LLM architecture achieves state-of-the-art results",
    "• IndicChat Pro: 85.2% accuracy across 11 languages + Hinglish",
    "• Offline-first design enables deployment in connectivity-limited areas",
    "• MuRIL BERT outperforms other models for Indian language tasks",
    "• Code-switching support addresses real-world usage patterns",
]
add_bullet_slide_content(slide, conclusions, Inches(1.0), Inches(2.95), Inches(5.1), Inches(3.5),
                          font_size=14, color=SOFT_WHITE)

# Right: Future Work
add_rounded_rect(slide, Inches(6.8), Inches(2.3), Inches(5.5), Inches(4.5), CARD_BG, ACCENT_GOLD)
add_text_box(slide, Inches(7.0), Inches(2.4), Inches(5.1), Inches(0.4),
             "🚀 Future Directions", font_size=20, color=ACCENT_GOLD, bold=True)

future = [
    "🔊 Voice-based interaction using fine-tuned Whisper ASR",
    "🌐 Expand to all 22 scheduled languages of India",
    "🔄 Advanced code-switching with neural LID models",
    "📱 Edge deployment on mobile devices",
    "🧠 Cross-lingual transfer learning for zero-shot support",
    "⚡ Energy-efficient architectures for scalable deployment",
    "🔒 Enhanced cultural sensitivity & bias mitigation",
]
add_bullet_slide_content(slide, future, Inches(7.0), Inches(2.95), Inches(5.1), Inches(3.5),
                          font_size=14, color=SOFT_WHITE)
add_slide_number(slide, 19, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
# SLIDE 20 — References
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_label(slide, "REFERENCES")
make_title(slide, "References")

refs = [
    '[1]  "Context-Aware Multilingual FAQ Chatbot Using MuRIL BERT" — Extractive QA with neural classifier',
    '[2]  "NLP-Based Chatbots for Hindi" — Healthcare & education chatbot development',
    '[3]  "ChatSense: Dynamic Knowledge Graphs" — Gemini 2.0 + real-time voice interaction',
    '[4]  "RAG with Knowledge Graph Integration" — Retrieval-Augmented Generation framework',
    '[5]  "English Bias in Large Language Models" — Analysis of LLM language representations',
    '[6]  "Whisper ASR Fine-Tuning for Indic Languages" — Joint multilingual training for WER reduction',
    '[7]  "Multi-Agent Hybrid Chatbot" — Generative + retrieval-based framework',
    '[8]  "Cross-Lingual Transfer Learning" — Adapter-based fine-tuning for low-resource languages',
    '[9]  "Chatbot Evaluation Frameworks" — Turing Test-based assessment',
    '[11] "IndicNLPSuite & IndicBERT" — IndicGLUE benchmark for Indian languages',
    '[12] "Intelligent Chatbots with ChatGPT" — LangChain + FAISS/Pinecone retrieval',
    '[13] "Survey of Chatbot Research Trends" — K-means clustering & ontology analysis',
]
add_bullet_slide_content(slide, refs, Inches(0.8), Inches(2.2), Inches(11.5), Inches(5.0),
                          font_size=13, color=SOFT_WHITE)
add_slide_number(slide, 20, TOTAL_SLIDES)


# ─── Save ─────────────────────────────────────────────────────────
output_path = os.path.join(r"e:\MAAM\chatbot code\docs",
                           "IndicChat_Pro_Presentation.pptx")
prs.save(output_path)
print(f"[OK] Presentation saved to: {output_path}")
print(f"     Total slides: {TOTAL_SLIDES}")
